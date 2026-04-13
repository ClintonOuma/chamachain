#!/usr/bin/env python3
"""
ChamaChain Presentation Generator
Creates a PowerPoint file from the presentation content.

Requirements: python-pptx library
Install: pip install python-pptx

Usage: python create_pptx.py
Output: ChamaChain_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Color scheme matching the web app
COLORS = {
    'background': RgbColor(10, 10, 15),      # #0a0a0f
    'primary': RgbColor(14, 165, 233),        # #0EA5E9
    'secondary': RgbColor(100, 116, 139),     # #64748B
    'text': RgbColor(248, 250, 252),          # #F8FAFC
    'success': RgbColor(134, 239, 172),       # #86EFAC
    'danger': RgbColor(252, 165, 165),        # #FCA5A5
    'accent': RgbColor(6, 182, 212),          # #06B6D4
}

def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, presenter_info):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # Presenter info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = presenter_info
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['text']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, is_two_column=False):
    """Create a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    if is_two_column and len(content_items) >= 2:
        # Two column layout
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        
        for item in content_items[0]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(12)
            p.level = 0
        
        right_box = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.2), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        
        for item in content_items[1]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(12)
            p.level = 0
    else:
        # Single column
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['text']
            p.space_before = Pt(14)
            p.level = 0
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Create a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(5)).table
    
    # Set column widths
    col_width = Inches(9 / num_cols)
    for col in table.columns:
        col.width = col_width
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(20, 30, 40)
    
    # Data rows
    for row_idx, row_data in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = COLORS['text']
            cell.text_frame.word_wrap = True
    
    return slide

def add_highlight_slide(prs, title, highlight_text, supporting_text=""):
    """Create a slide with a highlighted quote or statement."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Highlight box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2), Inches(8.4), Inches(3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(20, 35, 50)
    box.line.color.rgb = COLORS['primary']
    box.line.width = Pt(2)
    
    # Highlight text
    text_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(2.6))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = highlight_text
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = COLORS['text']
    p.alignment = PP_ALIGN.CENTER
    
    if supporting_text:
        support_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        tf = support_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = supporting_text
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_demo_slide(prs):
    """Create the demonstration overview slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demonstration"
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "5-Minute System Walkthrough"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    # Demo steps boxes
    steps = [
        ("1", "Google OAuth Login"),
        ("2", "Create Chama Group"),
        ("3", "AI Loan Application"),
        ("4", "Blockchain Proof")
    ]
    
    box_width = Inches(2)
    box_height = Inches(2)
    start_x = Inches(0.7)
    gap = Inches(0.4)
    y_pos = Inches(3)
    
    for i, (num, desc) in enumerate(steps):
        x_pos = start_x + (box_width + gap) * i
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(15, 30, 45)
        box.line.color.rgb = COLORS['primary']
        box.line.width = Pt(1.5)
        
        # Number
        num_box = slide.shapes.add_textbox(x_pos, y_pos + Inches(0.2), box_width, Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.1), y_pos + Inches(1), box_width - Inches(0.2), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER
    
    # Questions text
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you for your attention"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Generate the complete presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "ChamaChain",
        "AI-Powered Digital Chama Management Platform\nwith Blockchain Integration",
        "Presented by: [Your Name]\nInstitution: [Your University]\nDate: April 2026"
    )
    
    # Slide 2: Abstract - Two column
    add_content_slide(
        prs,
        "Abstract - Problem & Solution",
        [
            ["PROBLEM:",
             "• Manual record-keeping errors",
             "• Lack of transparency in funds",
             "• Limited credit for unbanked",
             "• No fraud detection",
             "• Inefficient coordination"],
            ["SOLUTION:",
             "• AI-powered credit scoring",
             "• Blockchain transparency",
             "• M-Pesa integration",
             "• Role-based access control",
             "• Real-time notifications"]
        ],
        is_two_column=True
    )
    
    # Metrics slide
    metrics_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(metrics_slide, COLORS['background'])
    
    title_box = metrics_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Achievements"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    metrics = [
        ("99.9%", "System Uptime"),
        ("<500ms", "AI Response Time"),
        ("40+", "API Endpoints")
    ]
    
    for i, (num, label) in enumerate(metrics):
        x = Inches(1 + i * 3.5)
        
        # Number
        num_box = metrics_slide.shapes.add_textbox(x, Inches(2.5), Inches(3), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = metrics_slide.shapes.add_textbox(x, Inches(3.5), Inches(3), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Objectives
    add_content_slide(
        prs,
        "Objectives",
        [
            "PRIMARY: Develop a scalable, secure, and intelligent digital platform that automates Chama operations while providing AI-driven financial insights to underbanked populations.",
            "",
            "SPECIFIC OBJECTIVES:",
            "• Digitize Chama operations (contributions, loans, meetings)",
            "• AI credit scoring using non-traditional data",
            "• Blockchain transparency for auditability",
            "• M-Pesa integration for mobile money",
            "• Fraud prevention with anomaly detection",
            "• Multi-tenancy with data isolation"
        ]
    )
    
    # Slide 4: Justification
    add_content_slide(
        prs,
        "Justification",
        [
            "ECONOMIC IMPACT:",
            "• Chamas contribute KES 400+ billion annually to Kenya's economy",
            "• 40% of Kenyan adults participate in rotating savings groups",
            "• Digital transformation reduces administrative costs by 60%",
            "• Addresses 4.4 million unbanked Kenyans in Chamas",
            "",
            "INNOVATION:",
            "• First platform combining AI + Blockchain + M-Pesa",
            "• Addresses UN SDG 1 (No Poverty) through financial inclusion",
            "• Novel credit scoring using group dynamics as alternative data"
        ]
    )
    
    # Slide 5: Limitations
    add_table_slide(
        prs,
        "Limitations",
        ["Constraint", "Impact", "Mitigation"],
        [
            ["AI Training Data", "6+ months needed", "Iterative improvement"],
            ["Blockchain Fees", "~$0.001/tx", "Polygon L2 (minimal)"],
            ["M-Pesa Limits", "KES 150K max/tx", "Batch processing"],
            ["Internet Required", "Needs connection", "PWA offline mode"],
            ["Regulatory", "CBK compliance", "Partner with lenders"]
        ]
    )
    
    # Slide 6: Context Diagram
    add_content_slide(
        prs,
        "System Context Diagram",
        [
            "EXTERNAL ENTITIES:",
            "• Chama Members (Admin, Treasurer, Member)",
            "• Super Administrators",
            "",
            "EXTERNAL SYSTEMS:",
            "• M-Pesa API - Mobile money payments",
            "• AI Service - Credit scoring & risk analysis",
            "• Google OAuth - Authentication",
            "• Polygon Blockchain - Immutable records",
            "• MongoDB Atlas - Data persistence"
        ]
    )
    
    # Slide 7: DFD
    add_content_slide(
        prs,
        "Data Flow: Loan Application Process",
        [
            "1. Member submits loan application (amount, purpose, duration)",
            "2. System fetches member history (contributions, existing loans)",
            "3. AI Service calculates credit score using alternative data",
            "4. Application saved to MongoDB with risk assessment",
            "5. Admin/Treasurer notified via push/email",
            "6. Approval/rejection sent to member with terms",
            "",
            "DATA STORES: User Database, Transaction History, Loan Applications, Blockchain Ledger"
        ]
    )
    
    # Slide 8: System Requirements
    add_table_slide(
        prs,
        "Technology Stack",
        ["Layer", "Technology", "Purpose"],
        [
            ["Frontend", "React 18 + Vite + Tailwind", "SPA, Responsive UI"],
            ["State", "Zustand", "Global state management"],
            ["Backend", "Node.js + Express", "REST API server"],
            ["Database", "MongoDB + Mongoose", "Document storage"],
            ["Auth", "JWT + bcrypt", "Token security"],
            ["AI", "Python + FastAPI", "ML credit scoring"],
            ["Blockchain", "Polygon (Mumbai)", "Smart contracts"],
            ["Payments", "M-Pesa Daraja API", "Mobile money"]
        ]
    )
    
    # Slide 9: Timeline
    add_content_slide(
        prs,
        "Project Timeline (6 Months)",
        [
            "Week 1-2:  Requirements Analysis - System analysis, user stories",
            "Week 3-4:  System Design - Database schema, UI mockups, API design",
            "Week 5-12: Core Development - Authentication, Chama CRUD, Members",
            "Week 13-16: AI Integration - Credit scoring model, risk assessment",
            "Week 17-20: Blockchain & M-Pesa - Smart contracts, mobile payments",
            "Week 21-24: Testing & Deployment - Unit tests, security audit, CI/CD",
            "",
            "ACTUAL HOURS: ~420 hours | CRITICAL PATH: Core Dev → AI → M-Pesa → Testing"
        ]
    )
    
    # Slide 10: Budget
    add_table_slide(
        prs,
        "Budget to Implement",
        ["Item", "Monthly Cost", "Notes"],
        [
            ["Render Backend", "FREE", "Starter plan sufficient"],
            ["Vercel Frontend", "FREE", "Free tier"],
            ["MongoDB Atlas M0", "FREE", "512MB cluster"],
            ["Cloudinary CDN", "FREE", "25GB free"],
            ["M-Pesa API Sandbox", "FREE", "Testing mode"],
            ["Google OAuth", "FREE", "<10K users"],
            ["Polygon Gas Fees", "~KES 50", "Minimal (~$0.35)"]
        ]
    )
    
    # Budget note
    budget_note = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(budget_note, COLORS['background'])
    
    note_box = budget_note.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Total Monthly Operating Cost: ~KES 50 (~$0.35)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['success']
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\nA traditional Chama secretary costs KES 5,000-10,000/month. This system serves unlimited members at a fraction of the cost."
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLORS['text']
    p2.alignment = PP_ALIGN.CENTER
    
    # Slide 11: Conclusion
    add_content_slide(
        prs,
        "Conclusion - Achievements",
        [
            "✅ Full-Stack Implementation: 40+ API endpoints, production-ready",
            "✅ AI Credit Scoring: Novel risk assessment using group behavior data",
            "✅ Blockchain Transparency: Immutable audit trail for all transactions",
            "✅ M-Pesa Integration: Seamless mobile money deposits and payouts",
            "✅ Security Architecture: JWT authentication, RBAC, rate limiting",
            "✅ 99.9% Uptime: Automated monitoring and cron jobs",
            "",
            "GitHub: github.com/ClintonOuma/chamachain",
            "Live Demo: chamachain-nine.vercel.app"
        ]
    )
    
    # Quote slide
    add_highlight_slide(
        prs,
        "",
        '"Financial inclusion starts with technology that understands local context. ChamaChain proves that AI and blockchain can serve the unbanked, not just the elite."',
        ""
    )
    
    # Slide 12: Recommendations
    add_content_slide(
        prs,
        "Recommendations & Future Work",
        [
            "IMMEDIATE (3 Months):",
            "• USSD Integration (*384#) for feature phone users",
            "• Mobile Apps (React Native) for iOS/Android",
            "• Biometric Authentication (fingerprint/face unlock)",
            "",
            "MEDIUM TERM (6-12 Months):",
            "• Credit Bureau Integration - Report to CRB for credit building",
            "• Insurance Module - Partner for micro-insurance products",
            "• Sacco Licensing - Apply for CBK regulatory compliance",
            "",
            "RESEARCH: Publish on 'Group Behavior as Alternative Credit Data'"
        ]
    )
    
    # Slide 13: Demo
    add_demo_slide(prs)
    
    # Save presentation
    output_path = "ChamaChain_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"⏱️  Estimated duration: 10 minutes + 5-minute demo")
    
    return prs

if __name__ == "__main__":
    try:
        from pptx import Presentation
        create_presentation()
    except ImportError:
        print("❌ python-pptx not installed.")
        print("Install with: pip install python-pptx")
        print("\nAlternative: Use the HTML slides (slides.html) which work in any browser!")
